"""
Phase 3, step 3: deterministic rendering of slide-content JSON into .pptx.

Deliberately separate from synthesis_agent.py: the LLM decides *what* the
story is (returns structured JSON), this module decides *how* it looks. This
split is what makes "auto-generate a 5-7 slide deck" reliable — no model is
ever asked to produce binary pptx content or fight with layout.

Palette: "Midnight Executive" (navy / ice blue / white) - a PMO / enterprise
delivery-governance register. Layout is data-driven off portfolio_package.json
(not just the LLM bullet text), so KPI counts, per-project score bars, and
status chips are always numerically consistent with the underlying JSON.
"""
import datetime as _dt

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY_DEEP = RGBColor(0x12, 0x18, 0x40)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x6F, 0x7A)
CARD_TINT = RGBColor(0xF3, 0xF5, 0xFB)
LINE_TINT = RGBColor(0xE1, 0xE5, 0xF0)

RAG_COLORS = {
    "Green": RGBColor(0x1E, 0x7D, 0x32),
    "Amber": RGBColor(0xC2, 0x7C, 0x0E),
    "Red": RGBColor(0xB0, 0x28, 0x28),
}
RAG_ORDER = ["Green", "Amber", "Red"]

# Severity (risk items) and priority (recommendation items) share the same
# three-tier color language as RAG status, so the deck reads consistently:
# red = urgent/critical, amber = elevated, navy = informational/lower.
SEVERITY_COLORS = {
    "Critical": RAG_COLORS["Red"],
    "High": RAG_COLORS["Amber"],
    "Medium": NAVY,
    "Low": MUTED,
}
PRIORITY_COLORS = {
    "P1": RAG_COLORS["Red"],
    "P2": RAG_COLORS["Amber"],
    "P3": NAVY,
    "P4": MUTED,
}

SEVERITY_COLORS = {
    "Critical": RGBColor(0x7A, 0x14, 0x14),
    "High": RAG_COLORS["Red"],
    "Medium": RAG_COLORS["Amber"],
    "Low": MUTED,
}
PRIORITY_COLORS = {
    "P1": RAG_COLORS["Red"],
    "P2": RAG_COLORS["Amber"],
    "P3": NAVY,
}
PRIORITY_LABELS = {
    "P1": "P1 \u2013 THIS WEEK",
    "P2": "P2 \u2013 NEXT REVIEW",
    "P3": "P3 \u2013 PROCESS FIX",
}


def _row_height(n_rows, top_in=1.75, bottom_in=7.02, gap_in=0.12, max_h=1.3, min_h=0.85):
    """Even row height that guarantees n_rows (+ gaps) fit between top_in and
    bottom_in — avoids the fixed-height overflow bug where a 5th/6th item
    would run under the footer."""
    n_rows = max(1, n_rows)
    available = bottom_in - top_in
    h = (available - (n_rows - 1) * gap_in) / n_rows
    return max(min_h, min(max_h, h))

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


# --------------------------------------------------------------------------
# low-level helpers
# --------------------------------------------------------------------------

def _blank_slide(prs, bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def _textbox(slide, text, left, top, width, height, size=16, color=DARK_TEXT,
             bold=False, italic=False, font="Calibri", align=PP_ALIGN.LEFT,
             anchor=None, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def _multiline(slide, lines, left, top, width, height, size=14, color=DARK_TEXT,
               bold=False, font="Calibri", align=PP_ALIGN.LEFT, space_after=6,
               bullet=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = (f"\u2022  {line}" if bullet else line)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(space_after)
    return tb


def _rect(slide, left, top, width, height, fill=None, line=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if radius:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _chip(slide, text, left, top, width, height, fill, text_color=WHITE, size=11):
    shp = _rect(slide, left, top, width, height, fill=fill, radius=0.5)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return shp


def _icon_circle(slide, glyph, left, top, diameter, fill, text_color=WHITE, size=16):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = glyph
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return shp


def _footer(slide, page_no, total_pages):
    _textbox(slide, "Project Health Reporting Agent  \u2013  Confidential",
             MARGIN, Inches(7.14), Inches(8.0), Inches(0.3),
             size=9, color=MUTED, font="Calibri")
    _textbox(slide, f"{page_no} / {total_pages}",
             Inches(12.2), Inches(7.14), Inches(0.55), Inches(0.3),
             size=9, color=MUTED, font="Calibri", align=PP_ALIGN.RIGHT)


def _header(slide, title, subtitle=None):
    _textbox(slide, title, MARGIN, Inches(0.42), Inches(12.1), Inches(0.6),
             size=28, color=NAVY, bold=True, font="Cambria")
    if subtitle:
        _textbox(slide, subtitle, MARGIN, Inches(0.98), Inches(12.1), Inches(0.4),
                 size=13, color=MUTED, font="Calibri")


# --------------------------------------------------------------------------
# slide builders
# --------------------------------------------------------------------------

def _add_title_slide(prs, slide_json, portfolio_package):
    slide = _blank_slide(prs, NAVY_DEEP)
    _rect(slide, Inches(0), Inches(5.55), SLIDE_W, Inches(1.95), fill=NAVY)

    _textbox(slide, "PMO EXECUTIVE SYNTHESIS", MARGIN, Inches(1.55), Inches(9), Inches(0.4),
             size=13, color=ICE, bold=True, font="Calibri")
    _textbox(slide, slide_json.get("title", "Portfolio Health Review"),
             MARGIN, Inches(2.05), Inches(12.0), Inches(1.5),
             size=42, color=WHITE, bold=True, font="Cambria")
    _textbox(slide, slide_json.get("subtitle", ""), MARGIN, Inches(3.35), Inches(11.5), Inches(0.6),
             size=17, color=ICE, font="Calibri")

    projects = portfolio_package.get("projects", [])
    band_mix = portfolio_package.get("band_mix", {})
    n = portfolio_package.get("generated_from_projects", len(projects))
    avg_score = round(sum(p.get("composite_score", 0) for p in projects) / n, 1) if n else 0
    red_n = band_mix.get("Red", 0)

    stats = [
        ("PROJECTS REVIEWED", str(n)),
        ("PORTFOLIO AVG SCORE", f"{avg_score}/100"),
        ("RED-FLAGGED", str(red_n)),
        ("REPORT DATE", _dt.date.today().strftime("%d %b %Y")),
    ]
    col_w = Inches(2.85)
    for i, (label, value) in enumerate(stats):
        x = MARGIN + i * (col_w + Inches(0.15))
        _textbox(slide, value, x, Inches(5.95), col_w, Inches(0.55),
                 size=26, color=WHITE, bold=True, font="Cambria")
        _textbox(slide, label, x, Inches(6.55), col_w, Inches(0.4),
                 size=10.5, color=ICE, font="Calibri")


def _add_rag_overview_slide(prs, slide_json, portfolio_package, page_no, total_pages):
    slide = _blank_slide(prs, WHITE)
    _header(slide, slide_json.get("title", "Portfolio RAG Overview"),
            "Composite status is computed from schedule, milestone-progress, blocker, and data-quality signals.")

    band_mix = portfolio_package.get("band_mix", {})
    projects = portfolio_package.get("projects", [])
    n = portfolio_package.get("generated_from_projects", len(projects))

    chart_data = CategoryChartData()
    chart_data.categories = RAG_ORDER
    chart_data.add_series("Projects", tuple(band_mix.get(b, 0) for b in RAG_ORDER))

    x, y, cx, cy = Inches(0.6), Inches(1.75), Inches(5.3), Inches(4.55)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data)
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)
    chart.legend.font.name = "Calibri"

    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0'
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(13)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = WHITE
    try:
        points = plot.series[0].points
        for point, label in zip(points, RAG_ORDER):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RAG_COLORS[label]
    except Exception:
        pass  # cosmetic only

    kpi = [
        ("TOTAL PROJECTS", str(n), NAVY),
        ("GREEN", str(band_mix.get("Green", 0)), RAG_COLORS["Green"]),
        ("AMBER", str(band_mix.get("Amber", 0)), RAG_COLORS["Amber"]),
        ("RED", str(band_mix.get("Red", 0)), RAG_COLORS["Red"]),
    ]
    card_w, card_h = Inches(2.95), Inches(1.35)
    gap = Inches(0.2)
    start_x, start_y = Inches(6.35), Inches(1.75)
    for i, (label, value, color) in enumerate(kpi):
        col, row = i % 2, i // 2
        cx0 = start_x + col * (card_w + gap)
        cy0 = start_y + row * (card_h + gap)
        _rect(slide, cx0, cy0, card_w, card_h, fill=CARD_TINT, radius=0.08)
        _rect(slide, cx0, cy0, Inches(0.08), card_h, fill=color)
        _textbox(slide, value, cx0 + Inches(0.25), cy0 + Inches(0.12), card_w - Inches(0.4), Inches(0.7),
                 size=32, color=color, bold=True, font="Cambria")
        _textbox(slide, label, cx0 + Inches(0.25), cy0 + Inches(0.9), card_w - Inches(0.4), Inches(0.35),
                 size=11, color=MUTED, bold=True, font="Calibri")

    note_y = start_y + 2 * (card_h + gap) + Inches(0.1)
    _multiline(slide, slide_json.get("bullets", []), start_x, note_y, Inches(6.1), Inches(1.0),
               size=12.5, color=MUTED, bullet=True, space_after=6)

    _footer(slide, page_no, total_pages)


def _add_project_snapshot_slide(prs, slide_json, portfolio_package, page_no, total_pages):
    slide = _blank_slide(prs, WHITE)
    _header(slide, slide_json.get("title", "Project-by-Project Snapshot"))

    projects = portfolio_package.get("projects", [])
    top = Inches(1.7)
    row_h = Inches(1.28)
    gap = Inches(0.14)
    card_w = Inches(12.13)

    max_rows = 4  # keep to one screen; overflow would move to an appendix slide in production
    for i, proj in enumerate(projects[:max_rows]):
        y = top + i * (row_h + gap)
        band = proj.get("current_rag", "Amber")
        color = RAG_COLORS.get(band, MUTED)
        score = proj.get("composite_score", 0)

        _rect(slide, MARGIN, y, card_w, row_h, fill=CARD_TINT, radius=0.06)
        _chip(slide, band.upper(), MARGIN + Inches(0.22), y + Inches(0.24), Inches(1.05), Inches(0.42), color)

        name = proj.get("project_name", "Untitled Project")
        driver = (proj.get("top_drivers") or ["No driver data available."])[0]
        if len(driver) > 108:
            driver = driver[:105] + "..."
        disagreement = proj.get("source_vs_computed_disagreement")

        _textbox(slide, name, MARGIN + Inches(1.5), y + Inches(0.15), Inches(7.7), Inches(0.4),
                 size=15, color=NAVY, bold=True, font="Calibri")
        _textbox(slide, driver, MARGIN + Inches(1.5), y + Inches(0.58), Inches(7.7), Inches(0.45),
                 size=11, color=MUTED, font="Calibri")
        if disagreement:
            _textbox(slide, f"\u26a0  PM-reported status: {disagreement.get('source_value')} \u2014 flagged for review",
                     MARGIN + Inches(1.5), y + Inches(0.94), Inches(7.7), Inches(0.3),
                     size=10, color=RAG_COLORS["Amber"], bold=True, font="Calibri")

        bar_x = MARGIN + Inches(9.4)
        _textbox(slide, f"{score:.0f}", bar_x, y + Inches(0.12), Inches(1.1), Inches(0.5),
                 size=24, color=color, bold=True, font="Cambria", align=PP_ALIGN.RIGHT)
        _textbox(slide, "/ 100", bar_x + Inches(1.05), y + Inches(0.27), Inches(0.6), Inches(0.3),
                 size=10, color=MUTED, font="Calibri")
        bar_w_total = Inches(2.3)
        _rect(slide, bar_x, y + Inches(0.72), bar_w_total, Inches(0.14), fill=LINE_TINT, radius=0.5)
        filled = Emu(int(bar_w_total * max(0.03, min(1.0, score / 100))))
        _rect(slide, bar_x, y + Inches(0.72), filled, Inches(0.14), fill=color, radius=0.5)

    _footer(slide, page_no, total_pages)


def _add_icon_row_slide(prs, slide_json, page_no, total_pages, icon_glyph="!", icon_color=NAVY,
                         numbered=False):
    """Fallback slide for plain bullet lists (legacy slide_plan.json shape,
    or any layout that hasn't been upgraded to structured items): icon-circle
    + bold lead-in + supporting text, one row per bullet."""
    slide = _blank_slide(prs, WHITE)
    _header(slide, slide_json.get("title", ""))

    bullets = slide_json.get("bullets", [])
    top = Inches(1.75)
    row_h = Inches(1.15)
    gap = Inches(0.12)

    if not bullets:
        _textbox(slide, "No items identified for this reporting period.", MARGIN, top, Inches(11.5), Inches(0.6),
                 size=14, color=MUTED, italic=True)
        _footer(slide, page_no, total_pages)
        return

    for i, bullet in enumerate(bullets[:5]):
        y = top + i * (row_h + gap)
        _rect(slide, MARGIN, y, Inches(12.13), row_h, fill=CARD_TINT, radius=0.08)
        glyph = str(i + 1) if numbered else icon_glyph
        _icon_circle(slide, glyph, MARGIN + Inches(0.22), y + Inches(0.28), Inches(0.6),
                     fill=icon_color, size=18)

        if ":" in bullet[:60]:
            lead, rest = bullet.split(":", 1)
            tb = slide.shapes.add_textbox(MARGIN + Inches(1.05), y + Inches(0.18), Inches(10.8), Inches(0.85))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = 0
            tf.margin_top = 0
            p = tf.paragraphs[0]
            r1 = p.add_run()
            r1.text = lead.strip() + ":  "
            r1.font.bold = True
            r1.font.size = Pt(14)
            r1.font.color.rgb = NAVY
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = rest.strip()
            r2.font.size = Pt(13)
            r2.font.color.rgb = DARK_TEXT
            r2.font.name = "Calibri"
        else:
            _textbox(slide, bullet, MARGIN + Inches(1.05), y + Inches(0.28), Inches(10.8), Inches(0.6),
                     size=13.5, color=DARK_TEXT, font="Calibri")

    _footer(slide, page_no, total_pages)


def _add_risk_items_slide(prs, slide_json, page_no, total_pages):
    slide = _blank_slide(prs, WHITE)
    _header(slide, slide_json.get("title", "Emerging Risks & Cross-Project Themes"),
            "Cross-project risk categories, ranked by severity \u2014 not a per-project status recap.")

    items = (slide_json.get("items") or [])[:5]
    top = Inches(1.75)
    gap = Inches(0.12)

    if not items:
        _textbox(slide, "No cross-project risks detected beyond current RAG status this period.",
                 MARGIN, top, Inches(11.5), Inches(0.6), size=14, color=MUTED, italic=True)
        _footer(slide, page_no, total_pages)
        return

    row_h = Inches(_row_height(len(items)))
    for i, item in enumerate(items):
        y = top + i * (row_h + gap)
        severity = item.get("severity", "Medium")
        color = SEVERITY_COLORS.get(severity, MUTED)

        _rect(slide, MARGIN, y, Inches(12.13), row_h, fill=CARD_TINT, radius=0.06)
        _chip(slide, severity.upper(), MARGIN + Inches(0.22), y + Inches(0.16), Inches(1.05), Inches(0.36),
              color, size=10)

        category = item.get("category", "Risk")
        statement = item.get("statement", "")
        affected = item.get("affected_projects") or []
        affected_tasks = item.get("affected_tasks") or []
        meta_bits = []
        if affected:
            meta_bits.append("Affects: " + ", ".join(affected))
        if affected_tasks:
            meta_bits.append("Tasks: " + ", ".join(affected_tasks[:3]))
        affected_line = "  \u00b7  ".join(meta_bits)

        _textbox(slide, category, MARGIN + Inches(1.5), y + Inches(0.1), Inches(10.4), Inches(0.32),
                 size=13.5, color=NAVY, bold=True, font="Calibri")
        _textbox(slide, statement, MARGIN + Inches(1.5), y + Inches(0.42), Inches(10.4),
                 row_h - Inches(0.65), size=11, color=DARK_TEXT, font="Calibri")
        if affected_line:
            av_top = y + row_h - Inches(0.28)
            _textbox(slide, affected_line, MARGIN + Inches(1.5), av_top, Inches(10.4), Inches(0.24),
                     size=9.5, color=MUTED, italic=True, font="Calibri")

    _footer(slide, page_no, total_pages)


def _add_recommendation_items_slide(prs, slide_json, page_no, total_pages):
    slide = _blank_slide(prs, WHITE)
    _header(slide, slide_json.get("title", "Recommendations & Decisions Needed"),
            "Prioritized by urgency \u2014 P1 acts this week, P3 is a systemic/process fix.")

    items = (slide_json.get("items") or [])[:5]
    top = Inches(1.75)
    gap = Inches(0.12)

    if not items:
        _textbox(slide, "No immediate actions required this period.", MARGIN, top, Inches(11.5), Inches(0.6),
                 size=14, color=MUTED, italic=True)
        _footer(slide, page_no, total_pages)
        return

    row_h = Inches(_row_height(len(items)))
    for i, item in enumerate(items):
        y = top + i * (row_h + gap)
        priority = item.get("priority", "P3")
        color = PRIORITY_COLORS.get(priority, NAVY)

        _rect(slide, MARGIN, y, Inches(12.13), row_h, fill=CARD_TINT, radius=0.06)
        _chip(slide, PRIORITY_LABELS.get(priority, priority), MARGIN + Inches(0.22), y + Inches(0.16),
              Inches(2.05), Inches(0.36), color, size=9.5)

        action = item.get("action", "")
        owner = item.get("owner", "")
        project = item.get("project", "")
        rationale = item.get("rationale", "")
        meta_line = "  \u00b7  ".join(x for x in [f"Owner: {owner}" if owner else "", project] if x)

        _textbox(slide, action, MARGIN + Inches(2.45), y + Inches(0.1), Inches(9.4), Inches(0.4),
                 size=13, color=NAVY, bold=True, font="Calibri")
        if meta_line:
            _textbox(slide, meta_line, MARGIN + Inches(2.45), y + Inches(0.47), Inches(9.4), Inches(0.24),
                     size=10, color=MUTED, bold=True, font="Calibri")
        if rationale:
            r_top = y + Inches(0.47) + (Inches(0.24) if meta_line else Inches(0))
            _textbox(slide, rationale, MARGIN + Inches(2.45), r_top, Inches(9.4),
                     row_h - (r_top - y) - Inches(0.05), size=10.5, color=DARK_TEXT, font="Calibri")

    _footer(slide, page_no, total_pages)


def _add_generic_bullets_slide(prs, slide_json, page_no, total_pages):
    slide = _blank_slide(prs, WHITE)
    _header(slide, slide_json.get("title", ""))
    _multiline(slide, slide_json.get("bullets", []), MARGIN, Inches(1.75), Inches(12.1), Inches(5.0),
               size=16, color=DARK_TEXT, bullet=True, space_after=12)
    _footer(slide, page_no, total_pages)


# --------------------------------------------------------------------------
# entry point
# --------------------------------------------------------------------------

def build_deck(slide_plan: dict, portfolio_package: dict, output_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = slide_plan.get("slides", [])
    total_content_pages = max(1, len([s for s in slides if s.get("layout") != "title"]))

    page_no = 0
    for slide_json in slides:
        layout = slide_json.get("layout")
        title_lower = (slide_json.get("title") or "").lower()

        if layout == "title":
            _add_title_slide(prs, slide_json, portfolio_package)
            continue

        page_no += 1
        if layout == "bullets_with_chart":
            _add_rag_overview_slide(prs, slide_json, portfolio_package, page_no, total_content_pages)
        elif layout == "risk_items":
            _add_risk_items_slide(prs, slide_json, page_no, total_content_pages)
        elif layout == "recommendation_items":
            _add_recommendation_items_slide(prs, slide_json, page_no, total_content_pages)
        elif "snapshot" in title_lower or "project-by-project" in title_lower:
            _add_project_snapshot_slide(prs, slide_json, portfolio_package, page_no, total_content_pages)
        elif "risk" in title_lower or "theme" in title_lower:
            # fallback path: older-shape slide_plan.json with plain bullets
            # instead of structured items (e.g. a real LLM call that ignored
            # the schema) - degrade to icon rows rather than failing to render.
            _add_icon_row_slide(prs, slide_json, page_no, total_content_pages,
                                 icon_glyph="!", icon_color=RAG_COLORS["Amber"])
        elif "recommend" in title_lower or "decision" in title_lower:
            _add_icon_row_slide(prs, slide_json, page_no, total_content_pages,
                                 icon_glyph="\u2713", icon_color=NAVY, numbered=True)
        else:
            _add_generic_bullets_slide(prs, slide_json, page_no, total_content_pages)

    prs.save(output_path)
    return output_path
